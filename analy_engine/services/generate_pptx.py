#!/usr/bin/env python3
"""
JSON → PPT 생성기
=================
디자인 원칙:
  - 전체 배경: 검정(#111111)
  - 기본 글씨: 흰색
  - 강조: 노란색(#FFD600)
  - 본문 폰트: 18pt (제시문·해설 모두 · 초과 시 슬라이드 분할)
  - 기본 정렬: 양쪽(JUSTIFY)
  - 원문 문단 구조(들여쓰기·문단 나눔) 보존
  - 제시문 끝에는 / 기호로 종료 표시
  - 누락/편집 없이 원문 그대로

슬라이드 순서 (문제별 묶음):
  [표지]
  문제 N마다:
    [섹션] "[문제 N]"
    [제시문] 18pt 기준 초과 시 자동 분할 · 문단 보존 · 끝에 /
    [문제(논제) + 요구사항/논제유형 표]
    [매트릭스 분석 — 강사 작성 영역]
    [매트릭스 예시답안]
    [대학 측 예시답안]  (univ_answer 필드)
    [출제의도]
    [문제해설]
    [채점기준 항목]
    [채점기준 등급표]
  [엔딩]
"""
import os, sys, json, re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("pip install python-pptx")
    sys.exit(1)

# ─────────────────────────────────────────────
# 상수
# ─────────────────────────────────────────────
C_BG    = RGBColor(0x11, 0x11, 0x11)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_YELL  = RGBColor(0xFF, 0xD6, 0x00)
C_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
C_DGRAY = RGBColor(0x44, 0x44, 0x44)
C_THEAD = RGBColor(0x2B, 0x2B, 0x2B)
C_TROW  = RGBColor(0x1A, 0x1A, 0x1A)
C_TROW2 = RGBColor(0x22, 0x22, 0x22)

GRADE_BG = {
    "상": RGBColor(0x10, 0x28, 0x10),
    "중": RGBColor(0x28, 0x20, 0x00),
    "하": RGBColor(0x28, 0x10, 0x10),
}

FONT     = "맑은 고딕"
SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)
L        = Inches(0.55)
CONT_W   = SLIDE_W - Inches(1.1)
T_HDR    = Inches(0.20)
T_BODY   = Inches(1.05)
BODY_H   = SLIDE_H - T_BODY - Inches(0.35)

# ── 본문 폰트 기준 ─────────────────────────────────────
BODY_PT         = 18     # 본문 기본 폰트 (제시문·해설 공통)
MIN_PT          = 16     # 표 안 등 보조 영역 최소 폰트
PARA_INDENT_PT  = 18     # 문단 첫 글자 들여쓰기 (전각 공백 1칸 상당)
PARA_SPACE_AFTER = 6     # 문단 간 간격(pt) — 문단 구분감 확보
LINE_SPACE_AFTER = 2     # 같은 문단 내 줄 간격(pt)

# ── 한 슬라이드 본문 용량 (18pt + 줄간격 1.4 기준) ─────
# 슬라이드 본문 영역: 12.2인치(가로) × 약 6.1인치(세로)
# 18pt 맑은고딕 한글: 한 줄 약 38~42자
# 줄간격 1.4 적용 시 한 페이지 약 17~18줄
# → 문단 공백·들여쓰기 감안 안전치 680자
CHARS_PER_SLIDE = 680

# ── 줄간격 (배수) ─────────────────────────────────────
# 1.0 = 빡빡함(기본), 1.4 = 가독성 좋음, 1.5 = 넉넉함
LINE_SPACING = 1.4


# ─────────────────────────────────────────────
# 기본 헬퍼
# ─────────────────────────────────────────────
def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = C_BG


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    return box.text_frame


def p(tf, text, size=BODY_PT, bold=False, color=None,
      align=PP_ALIGN.JUSTIFY, sb=0, sa=3, first=False,
      indent_pt=0, line_spacing=None):
    """
    단락 추가.
    indent_pt: 문단 첫 줄 들여쓰기 (pt 단위 · 0이면 없음)
    line_spacing: 줄간격 배수 (None이면 PPT 기본값)
    """
    color = color or C_WHITE
    if first and tf.paragraphs and tf.paragraphs[0].text == "":
        par = tf.paragraphs[0]
    else:
        par = tf.add_paragraph()
    par.alignment    = align
    par.space_before = Pt(sb)
    par.space_after  = Pt(sa)
    if line_spacing is not None:
        par.line_spacing = line_spacing
    # 문단 첫 줄 들여쓰기 (python-pptx 직접 API 없어서 XML 설정)
    if indent_pt > 0:
        pPr = par._pPr if par._pPr is not None else par._p.get_or_add_pPr()
        # marL=0, indent=양수 값 → 첫 줄 들여쓰기 (EMU 단위: 1pt = 12700 EMU)
        pPr.set('marL', '0')
        pPr.set('indent', str(int(indent_pt * 12700)))
    run = par.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color
    return par


def hline(slide, top, color=C_DGRAY):
    s = slide.shapes.add_shape(1, L, top, CONT_W, Emu(14000))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def header(slide, title, sub=""):
    tf_ = tb(slide, L, T_HDR, CONT_W, Inches(0.6))
    p(tf_, title, size=19, bold=True, color=C_YELL,
      align=PP_ALIGN.LEFT, first=True)
    if sub:
        p(tf_, f"  {sub}", size=16, color=C_GRAY, align=PP_ALIGN.LEFT)
    hline(slide, Inches(0.90))


# ─────────────────────────────────────────────
# 문단 파싱 & 분할 (신규)
# ─────────────────────────────────────────────
def split_into_paragraphs(text):
    """
    원문 텍스트를 문단 리스트로 분할.
    구분 기준:
      - 빈 줄 (\n\n 이상) → 확실한 문단 경계
      - 단일 줄바꿈 (\n) → 같은 문단 내 줄바꿈으로 간주하고 공백으로 이어붙임
    반환: [문단1, 문단2, ...]
    """
    if not text:
        return []
    # 빈 줄(연속 \n) 기준으로 문단 분리
    raw_paras = re.split(r'\n\s*\n+', text.strip())
    result = []
    for rp in raw_paras:
        # 같은 문단 내 줄바꿈은 공백으로 이어 붙이되, 연속 공백은 정리
        joined = re.sub(r'\s*\n\s*', ' ', rp.strip())
        joined = re.sub(r' {2,}', ' ', joined)
        if joined:
            result.append(joined)
    return result if result else [text.strip()]


def split_paragraphs_by_capacity(paragraphs, cap=CHARS_PER_SLIDE):
    """
    문단 리스트를 슬라이드 용량에 맞춰 청크로 분할.
    각 청크도 문단 리스트 형태.

    원칙:
      1. 문단 경계를 최대한 존중 (한 문단이 두 슬라이드에 걸치지 않도록 우선)
      2. 단, 한 문단이 cap 초과 시 어쩔 수 없이 문장 단위로 분할
    반환: [[문단1, 문단2], [문단3], [문단4, 문단5, 문단6], ...]
    """
    chunks = []
    cur = []
    cur_len = 0

    for para in paragraphs:
        plen = len(para)

        # 단일 문단이 cap 초과 → 문장 단위 분할
        if plen > cap:
            # 현재 쌓인 것 먼저 flush
            if cur:
                chunks.append(cur)
                cur, cur_len = [], 0
            # 긴 문단을 문장 단위로 분할
            sentences = _split_sentences(para)
            sub_cur = []
            sub_len = 0
            for sent in sentences:
                slen = len(sent)
                if sub_cur and sub_len + slen > cap:
                    chunks.append([' '.join(sub_cur)])
                    sub_cur, sub_len = [], 0
                sub_cur.append(sent)
                sub_len += slen + 1
            if sub_cur:
                chunks.append([' '.join(sub_cur)])
            continue

        # 현재 청크에 추가하면 초과 → 새 청크 시작
        if cur and cur_len + plen > cap:
            chunks.append(cur)
            cur, cur_len = [], 0

        cur.append(para)
        cur_len += plen + 2  # 문단 간 공백 가중치

    if cur:
        chunks.append(cur)
    return chunks if chunks else [paragraphs]


def _split_sentences(text):
    """한국어 문장 단위 분할 (마침표·물음표·느낌표 + 종결어미)"""
    # 종결 지점 뒤에 공백이 오면 문장 경계로 간주
    sents = re.split(r'(?<=[.!?。！？])\s+', text)
    return [s.strip() for s in sents if s.strip()]


def tbl(slide, rows, cols, l, t, w, h):
    return slide.shapes.add_table(rows, cols, l, t, w, h).table


def cell(c, text, size=MIN_PT, bold=False, color=None,
         align=PP_ALIGN.LEFT, bg_c=None):
    color = color or C_WHITE
    if bg_c:
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
    tf_ = c.text_frame
    tf_.word_wrap = True
    par = tf_.paragraphs[0]
    par.alignment = align
    run = par.add_run()
    run.text           = str(text)
    run.font.name      = FONT
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color


# ─────────────────────────────────────────────
# 본문 렌더링 헬퍼 (신규) — 문단 구조 보존
# ─────────────────────────────────────────────
def render_body_paragraphs(tf, paragraphs, size=BODY_PT,
                            indent=True, end_slash=False,
                            end_tilde=False):
    """
    문단 리스트를 텍스트 프레임에 렌더링.
      - 첫 글자 들여쓰기 (indent=True)
      - 문단 간 공백 (PARA_SPACE_AFTER)
      - 줄간격 (LINE_SPACING 배수)
      - 양쪽 정렬
      - end_slash=True 이면 마지막 문단 끝에 ' /' 추가 (제시문 종료 표시)
      - end_tilde=True 이면 마지막 문단 끝에 ' ~' 추가 (다음 슬라이드 이어짐 표시)
    """
    first = True
    n = len(paragraphs)
    for idx, para in enumerate(paragraphs):
        text = para
        is_last_para = (idx == n - 1)
        if is_last_para:
            if end_slash:
                text = text.rstrip() + "  /"
            elif end_tilde:
                text = text.rstrip() + "  ~"
        p(tf, text,
          size=size,
          color=C_WHITE,
          align=PP_ALIGN.JUSTIFY,
          sa=PARA_SPACE_AFTER,
          first=first,
          indent_pt=PARA_INDENT_PT if indent else 0,
          line_spacing=LINE_SPACING)
        first = False


# ─────────────────────────────────────────────
# 슬라이드 빌더
# ─────────────────────────────────────────────

def title_slide(prs, meta):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    univ = meta.get("university", "")
    year = meta.get("year", "")
    track = meta.get("track", "")
    sub = meta.get("subtitle", "")
    t = meta.get("examTime", "")

    tf_ = tb(sl, L, Inches(1.5), CONT_W, Inches(1.7))
    p(tf_, f"{univ} 파이널", 54, True, C_WHITE, PP_ALIGN.CENTER, first=True)

    tf_ = tb(sl, L, Inches(3.3), CONT_W, Inches(0.9))
    p(tf_, f"{year} 기출  ({track})", 30, True, C_YELL, PP_ALIGN.CENTER, first=True)

    if sub:
        tf_ = tb(sl, L, Inches(4.3), CONT_W, Inches(0.65))
        p(tf_, f"- {sub} -", 22, False, C_GRAY, PP_ALIGN.CENTER, first=True)

    if t:
        tf_ = tb(sl, L, Inches(5.1), CONT_W, Inches(0.5))
        p(tf_, f"시험 시간 {t}분", 16, False, C_DGRAY, PP_ALIGN.CENTER, first=True)

    tf_ = tb(sl, L, Inches(6.6), CONT_W, Inches(0.5))
    p(tf_, "EOLE 논술 연구소  |  www.eole.co.kr", 13, False,
      C_DGRAY, PP_ALIGN.CENTER, first=True)


def section_slide(prs, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tf_ = tb(sl, L, Inches(2.5), CONT_W, Inches(1.6))
    p(tf_, f"[ 문제  {num} ]", 60, True, C_YELL, PP_ALIGN.CENTER, first=True)
    tf_ = tb(sl, L, Inches(4.3), CONT_W, Inches(0.8))
    p(tf_, "제시문 분석,  매트릭스와 예시답안", 22, False,
      C_GRAY, PP_ALIGN.CENTER, first=True)


def passage_slide(prs, label, paragraphs, source="", textbook="",
                  pg="", is_last=False):
    """
    제시문 슬라이드.
    paragraphs: 이 슬라이드에 들어갈 문단 리스트
    is_last: 해당 제시문의 마지막 청크 여부
             - True  → 본문 끝에 ' /' + 출처 표시 (제시문 종료)
             - False → 본문 끝에 ' ~' 표시 (다음 슬라이드로 이어짐)
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, label + (f"  {pg}" if pg else ""))

    tf_ = tb(sl, L, T_BODY, CONT_W, BODY_H - Inches(0.4))
    render_body_paragraphs(tf_, paragraphs,
                            size=BODY_PT,
                            indent=True,
                            end_slash=is_last,
                            end_tilde=(not is_last))

    # 출처는 마지막 청크에만
    if is_last:
        parts = []
        if source:   parts.append(source)
        if textbook: parts.append("\u300e" + textbook.strip("\u300e\u300f") + "\u300f")
        if parts:
            tf_ = tb(sl, L, SLIDE_H - Inches(0.5), CONT_W, Inches(0.38))
            p(tf_, "  /  ".join(parts), 14, False,
              C_GRAY, PP_ALIGN.RIGHT, first=True)


def question_slide(prs, num, question, instructions=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, f"[문제 {num}]  논제")

    q_text = question.get("text", "")
    wc     = question.get("wordCount", "")
    pts    = question.get("points", "")

    y = T_BODY
    if instructions:
        tf_ = tb(sl, L, y, CONT_W, Inches(0.42))
        p(tf_, instructions, 16, False, C_GRAY, PP_ALIGN.LEFT, first=True)
        y += Inches(0.45)

    # 논제 본문: 18pt + 양쪽 정렬 + 들여쓰기 없음(논제라 짧음)
    tf_ = tb(sl, L, y, CONT_W, Inches(2.55))
    q_paragraphs = split_into_paragraphs(q_text)
    render_body_paragraphs(tf_, q_paragraphs,
                            size=BODY_PT, indent=False, end_slash=False)
    if wc or pts:
        parts = []
        if wc:  parts.append(f"<{wc}>")
        if pts: parts.append(f"[{pts}점]")
        p(tf_, "  ".join(parts), 16, True, C_YELL, PP_ALIGN.RIGHT, sb=4)

    sep_y = y + Inches(2.8)
    hline(sl, sep_y)

    tbl_y = sep_y + Inches(0.12)
    tbl_h = SLIDE_H - tbl_y - Inches(0.2)
    t_ = tbl(sl, 2, 2, L, tbl_y, CONT_W, tbl_h)
    t_.columns[0].width = Inches(1.6)
    t_.columns[1].width = CONT_W - Inches(1.6)

    cell(t_.cell(0, 0), "요구사항", 16, True, C_YELL, PP_ALIGN.CENTER, C_THEAD)
    nums = re.findall(r"[①②③④⑤⑥][^①②③④⑤⑥]*", q_text)
    req = "\n".join(f"{i+1}) {n.strip()}" for i, n in enumerate(nums)) if nums else q_text
    cell(t_.cell(0, 1), req, 16, False, C_WHITE, PP_ALIGN.LEFT, C_TROW)

    cell(t_.cell(1, 0), "논제유형", 16, True, C_YELL, PP_ALIGN.CENTER, C_THEAD)
    types = []
    if re.search(r"분류|나누|구분",   q_text): types.append("분류")
    if re.search(r"요약|정리",        q_text): types.append("요약")
    if re.search(r"설명|서술",        q_text): types.append("설명")
    if re.search(r"비판|비교",        q_text): types.append("비판/비교")
    if re.search(r"근거|정당화|옹호", q_text): types.append("논증")
    if re.search(r"활용|자료",        q_text): types.append("자료해석")
    cell(t_.cell(1, 1), ", ".join(types) if types else "서술",
         14, False, C_WHITE, PP_ALIGN.LEFT, C_TROW2)


def matrix_slide(prs, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, f"◎ 매트릭스 분석  [문제 {num}]", "강사 작성 영역")

    labels = ["기준", "대상", "지점", "내용"]
    t_ = tbl(sl, len(labels), 2, L, T_BODY, CONT_W, BODY_H)
    t_.columns[0].width = Inches(1.5)
    t_.columns[1].width = CONT_W - Inches(1.5)
    for i, lbl in enumerate(labels):
        cell(t_.cell(i, 0), lbl, 16, True, C_YELL, PP_ALIGN.CENTER, C_THEAD)
        cell(t_.cell(i, 1), "", 16, False, C_WHITE,
             bg_c=C_TROW if i % 2 == 0 else C_TROW2)


def text_slides(prs, hdr, text, sub=""):
    """
    범용 텍스트 슬라이드 — 18pt · 문단 보존 · 용량 초과 시 분할
    (해설·예시답안·채점기준 공용)
    """
    paragraphs = split_into_paragraphs(text)
    chunks = split_paragraphs_by_capacity(paragraphs, CHARS_PER_SLIDE)
    total  = len(chunks)
    for ci, chunk_paras in enumerate(chunks):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg(sl)
        pg = f"({ci+1}/{total})" if total > 1 else ""
        header(sl, hdr + (f"  {pg}" if pg else ""), sub)
        tf_ = tb(sl, L, T_BODY, CONT_W, BODY_H)
        render_body_paragraphs(tf_, chunk_paras,
                                size=BODY_PT,
                                indent=True,
                                end_slash=False)


def rubric_table_slide(prs, num, rubric_table):
    if not rubric_table:
        return
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    header(sl, f"◎ 대학측 채점 기준  [문제 {num}]", "등급표")

    nr = len(rubric_table) + 1
    rh = min(BODY_H / nr, Inches(0.88))
    t_ = tbl(sl, nr, 3, L, T_BODY, CONT_W, rh * nr)
    t_.columns[0].width = Inches(1.2)
    t_.columns[1].width = Inches(1.0)
    t_.columns[2].width = CONT_W - Inches(2.2)

    for j, h_ in enumerate(["등급", "코드", "기준"]):
        cell(t_.cell(0, j), h_, 16, True, C_YELL, PP_ALIGN.CENTER, C_THEAD)

    for i, item in enumerate(rubric_table):
        ri    = i + 1
        grade = item.get("grade", "")
        code  = item.get("code", "")
        desc  = item.get("desc", "")
        gb    = GRADE_BG.get(grade, C_TROW)
        fs    = MIN_PT if len(desc) < 55 else 14
        cell(t_.cell(ri, 0), grade, 16, True,  C_YELL,  PP_ALIGN.CENTER, gb)
        cell(t_.cell(ri, 1), code,  16, True,  C_WHITE, PP_ALIGN.CENTER, gb)
        cell(t_.cell(ri, 2), desc,  fs, False, C_WHITE, PP_ALIGN.LEFT,   gb)


def ending_slide(prs, meta):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tf_ = tb(sl, L, Inches(2.6), CONT_W, Inches(1.2))
    p(tf_, f"{meta.get('university','')} 파이널", 44, True,
      C_WHITE, PP_ALIGN.CENTER, first=True)
    tf_ = tb(sl, L, Inches(3.9), CONT_W, Inches(0.8))
    p(tf_, f"{meta.get('year','')} 기출  —  끝  —", 24, False,
      C_YELL, PP_ALIGN.CENTER, first=True)
    tf_ = tb(sl, L, Inches(6.6), CONT_W, Inches(0.5))
    p(tf_, "EOLE 논술 연구소  |  www.eole.co.kr", 13, False,
      C_DGRAY, PP_ALIGN.CENTER, first=True)


# ─────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────
def generate_pptx(data: dict, output_path: str):
    meta = data.get("meta", {})
    sets = data.get("problemSets", [])

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, meta)

    for ps in sets:
        num          = ps.get("number", 1)
        instructions = ps.get("instructions", "")
        question     = ps.get("question", {})

        section_slide(prs, num)

        # 제시문 분할 (문단 기반)
        for psg in ps.get("passages", []):
            label    = psg.get("label", "")
            text     = psg.get("text", "")
            source   = psg.get("source", "")
            textbook = psg.get("textbook", "")

            # 1) 원문 → 문단 리스트
            paragraphs = split_into_paragraphs(text)
            # 2) 문단 리스트 → 용량 기준 청크 분할
            chunks = split_paragraphs_by_capacity(paragraphs, CHARS_PER_SLIDE)
            total = len(chunks)

            for ci, chunk_paras in enumerate(chunks):
                pg      = f"({ci+1}/{total})" if total > 1 else ""
                is_last = ci == total - 1
                passage_slide(prs, label, chunk_paras,
                              source, textbook, pg, is_last)

        # 문제
        if question:
            question_slide(prs, num, question, instructions)

        # 매트릭스
        matrix_slide(prs, num)

        # 예시답안들
        if ps.get("sampleAnswer"):
            text_slides(prs, f"[매트릭스 예시답안]  [문제 {num}]",
                        ps["sampleAnswer"])
        if ps.get("univ_answer"):
            text_slides(prs, f"[대학 측 예시답안]  [문제 {num}]",
                        ps["univ_answer"])

        # 해설
        if ps.get("출제의도"):
            text_slides(prs, f"◎ 출제의도  [문제 {num}]", ps["출제의도"])
        if ps.get("문제해설"):
            text_slides(prs, f"◎ 문제해설  [문제 {num}]", ps["문제해설"])

        # 채점기준
        if ps.get("rubric"):
            text_slides(prs, f"◎ 대학측 채점 기준  [문제 {num}]", ps["rubric"])
        if ps.get("rubricTable"):
            rubric_table_slide(prs, num, ps["rubricTable"])

    ending_slide(prs, meta)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    kb = os.path.getsize(output_path) // 1024
    print(f"✅ PPT 생성 완료: {output_path}  ({kb:,} KB, {len(prs.slides)}슬라이드)")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python generate_pptx.py <data.json> [output.pptx]")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        data = json.load(f)
    meta = data.get("meta", {})
    default = f"{meta.get('university','')}_{meta.get('year','')}_{meta.get('track','').replace(' ','')}.pptx"
    output = sys.argv[2] if len(sys.argv) >= 3 else f"outputs/{default}"
    generate_pptx(data, output)
